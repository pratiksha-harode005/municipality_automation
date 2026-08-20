import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 Widescreen dimensions
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Clean, harmonious color palette
    COLOR_NAVY = RGBColor(11, 47, 69)        # #0B2F45 (Deep Royal Navy)
    COLOR_CYAN = RGBColor(0, 139, 149)       # #008B95 (Municipal Teal/Cyan)
    COLOR_LIGHT_CYAN = RGBColor(56, 189, 248)# #38BDF8 (Sky Blue)
    COLOR_GOLD = RGBColor(217, 119, 6)       # #D97706 (Amber Gold)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_CARD_BG = RGBColor(255, 255, 255)
    COLOR_CARD_BORDER = RGBColor(226, 232, 240)
    COLOR_TEXT_MAIN = RGBColor(30, 41, 59)   # Slate 800
    COLOR_TEXT_MUTED = RGBColor(100, 116, 139)
    COLOR_SUCCESS_BG = RGBColor(240, 253, 244)
    COLOR_SUCCESS_BORDER = RGBColor(187, 247, 208)
    COLOR_ALERT_BG = RGBColor(254, 242, 242)
    COLOR_ALERT_BORDER = RGBColor(254, 202, 202)

    def add_header(slide, title_text, category_text="MUNICIPALITY AUTOMATION SYSTEM"):
        # Header category pill
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.32))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_CYAN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.733), Inches(0.65))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(23)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_NAVY

        # Subtle decorative divider line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.42), Inches(11.733), Inches(0.02)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_CYAN
        line.line.color.rgb = COLOR_CYAN

    def add_card(slide, left, top, width, height, title, items, badge=None, bg_color=COLOR_WHITE, border_color=COLOR_CARD_BORDER, title_color=COLOR_NAVY):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.24)
        tf.margin_right = Inches(0.24)
        tf.margin_top = Inches(0.22)
        tf.margin_bottom = Inches(0.22)

        if badge:
            p_badge = tf.paragraphs[0]
            p_badge.text = badge.upper()
            p_badge.font.size = Pt(9.5)
            p_badge.font.bold = True
            p_badge.font.color.rgb = COLOR_CYAN
            p_title = tf.add_paragraph()
        else:
            p_title = tf.paragraphs[0]

        p_title.text = title
        p_title.font.size = Pt(15)
        p_title.font.bold = True
        p_title.font.color.rgb = title_color
        p_title.space_after = Pt(8)

        for item in items:
            p_item = tf.add_paragraph()
            p_item.text = f"•  {item}"
            p_item.font.size = Pt(11.5)
            p_item.font.color.rgb = COLOR_TEXT_MAIN if bg_color == COLOR_WHITE or bg_color == COLOR_SUCCESS_BG or bg_color == COLOR_ALERT_BG else COLOR_WHITE
            p_item.space_after = Pt(5)

        return card

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_NAVY
    bg1.line.fill.background()

    strip = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), Inches(7.5))
    strip.fill.solid()
    strip.fill.fore_color.rgb = COLOR_CYAN
    strip.line.fill.background()

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(1.3), Inches(11.0), Inches(4.8))
    tf1 = tb1.text_frame
    tf1.word_wrap = True

    p1 = tf1.paragraphs[0]
    p1.text = "🏛️ MUNICIPAL CORPORATION AUTOMATION SYSTEM"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_LIGHT_CYAN
    p1.space_after = Pt(14)

    p2 = tf1.add_paragraph()
    p2.text = "Digital E-Governance & Civic Services Platform"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.space_after = Pt(16)

    p3 = tf1.add_paragraph()
    p3.text = "A modern, complete digital solution designed to make municipal services fast, transparent, and accessible for Citizens, Officers, and City Administrators."
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_after = Pt(32)

    p4 = tf1.add_paragraph()
    p4.text = "⚡ Key Highlights:  100% Online  |  Photo Evidence  |  Live Tracking  |  Mobile & Laptop Friendly"
    p4.font.size = Pt(12.5)
    p4.font.bold = True
    p4.font.color.rgb = COLOR_GOLD

    # ==========================================
    # SLIDE 2: WHAT IS THIS PROJECT?
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, "What is This Project? (Simple Overview)", "PROJECT INTRODUCTION")

    add_card(s2, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "What is this Platform?",
             [
                 "A complete digital website built for Municipal Corporations (like BBMP).",
                 "Connects citizens directly with city government departments online.",
                 "Citizens can apply for civic services and report local city issues from home.",
                 "Municipal officers can review, inspect, and resolve issues digitally.",
                 "Eliminates paperwork and saves people hours of waiting in government offices."
             ],
             badge="SIMPLE DEFINITION")

    add_card(s2, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "Who is it Built For?",
             [
                 "Citizens: File complaints, apply for permits, and download certificates.",
                 "Municipal Officers: Manage assigned ward tasks and conduct inspections.",
                 "City Administrators: Oversee all wards, departments, staff, and city notices.",
                 "Public Visitors: Explore city directory, public hearing notices, and interactive maps."
             ],
             badge="AUDIENCE & USERS")

    # ==========================================
    # SLIDE 3: PROBLEMS WE ARE SOLVING
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, "Real-World Problems We Are Solving", "THE CHALLENGE")

    add_card(s3, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "Old / Traditional System Issues",
             [
                 "Long Queues: Citizens had to stand in lines for simple forms and certificates.",
                 "No Tracking: After submitting an application, nobody knew where it was stuck.",
                 "Unreported Damage: Potholes, broken streetlights, and garbage piles took weeks to fix.",
                 "Lost Paperwork: Physical paper applications were frequently misplaced.",
                 "Lack of Transparency: Citizens had no direct contact with responsible ward officers."
             ],
             badge="BEFORE AUTOMATION",
             bg_color=COLOR_ALERT_BG,
             border_color=COLOR_ALERT_BORDER)

    add_card(s3, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "How Our System Fixes Them",
             [
                 "Apply from Anywhere: Submit requests in 2 minutes using mobile phone or PC.",
                 "Live Tracking ID: Watch your application move from 'Submitted' to 'Resolved'.",
                 "Photo Proof: Click and upload live photos of road issues or garbage blackspots.",
                 "Safe Cloud Storage: All applications are stored permanently in a secure database.",
                 "Direct Accountability: Every action is logged with the officer's name and timestamp."
             ],
             badge="AFTER AUTOMATION",
             bg_color=COLOR_SUCCESS_BG,
             border_color=COLOR_SUCCESS_BORDER)

    # ==========================================
    # SLIDE 4: THE 3 MAIN USER ROLES
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, "The 3 User Roles in the System", "USER MANAGEMENT")

    add_card(s4, Inches(0.8), Inches(1.75), Inches(3.6), Inches(5.1),
             "1. The Citizen",
             [
                 "Creates personal account in seconds.",
                 "Applies for road repairs, birth/death records, and water connections.",
                 "Uploads photo evidence.",
                 "Tracks live progress status.",
                 "Downloads official approved receipts."
             ],
             badge="CITIZEN PORTAL")

    add_card(s4, Inches(4.85), Inches(1.75), Inches(3.6), Inches(5.1),
             "2. The Municipal Officer",
             [
                 "Dedicated officer work dashboard.",
                 "Views all complaints in their ward.",
                 "Inspects photo proof and GPS location.",
                 "Takes action and adds resolution notes.",
                 "Marks complaints as Resolved."
             ],
             badge="OFFICER SUITE")

    add_card(s4, Inches(8.9), Inches(1.75), Inches(3.6), Inches(5.1),
             "3. Super Administrator",
             [
                 "Master control over entire portal.",
                 "Manages officer accounts & permissions.",
                 "Updates city corporations and ward info.",
                 "Publishes emergency notices and alerts.",
                 "Views real-time analytics & reports."
             ],
             badge="ADMIN CMS")

    # ==========================================
    # SLIDE 5: MAJOR CIVIC SERVICES (PART 1)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, "Major Civic Services: Roads & Vital Records", "SERVICES OVERVIEW")

    add_card(s5, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "🛣️ Road, Pothole & Streetlight Reporting",
             [
                 "FixMyCity Feature: Report road damage, potholes, or dark streetlights.",
                 "Visual Photo Upload: Take a photo directly from your phone camera.",
                 "Landmark Address: Enter street name, ward number, and GPS landmark.",
                 "Fast Response: Assigned directly to the ward civil engineering team.",
                 "Real-time Notification: Get notified when the road is repaired."
             ],
             badge="CIVIL ENGINEERING")

    add_card(s5, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "📜 Birth & Death Certificate Registration",
             [
                 "Digital Application: Register births or deaths without visiting municipal office.",
                 "Document Upload: Attach hospital discharge summary or identity proof.",
                 "Digital Verification: Medical health officer reviews and approves online.",
                 "Instant Certificate: Download verified digital certificate with official QR code."
             ],
             badge="HEALTH & VITAL RECORDS")

    # ==========================================
    # SLIDE 6: MAJOR CIVIC SERVICES (PART 2)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, "Major Civic Services: Water & Waste Management", "SERVICES OVERVIEW")

    add_card(s6, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "💧 Water Connection & Drainage Permitting",
             [
                 "New Connection Requests: Apply for home or commercial water pipelines.",
                 "Property Khata Verification: Enter property details for automated verification.",
                 "Inspection Booking: Schedule a site inspection by water board engineers.",
                 "Transparent Pricing: Clear estimation of connection fees and water charges."
             ],
             badge="WATER UTILITY")

    add_card(s6, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "🚚 Solid Waste & Tipper Timetable",
             [
                 "Ward Tipper Timetable: Check daily garbage truck arrival times in your area.",
                 "Waste Segregation Guide: Easy tips for wet, dry, and sanitary waste separation.",
                 "Report Garbage Blackspots: Upload photos of uncollected garbage for rapid cleaning.",
                 "Cleanliness Drive Tracking: Updates on local sanitation and park cleanup drives."
             ],
             badge="SANITATION & CLEANLINESS")

    # ==========================================
    # SLIDE 7: STEP-BY-STEP CITIZEN JOURNEY
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, "Step-by-Step Citizen Journey (How it Works)", "CITIZEN WORKFLOW")

    add_card(s7, Inches(0.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Step 1: Apply",
             [
                 "Log in to the portal.",
                 "Choose your civic service (e.g. Pothole repair).",
                 "Fill in simple details in under 2 minutes."
             ],
             badge="SUBMISSION")

    add_card(s7, Inches(3.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Step 2: Attach Evidence",
             [
                 "Upload clear photos from phone or gallery.",
                 "Enter landmark or ward location.",
                 "Submit your request."
             ],
             badge="EVIDENCE")

    add_card(s7, Inches(6.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Step 3: Track Live",
             [
                 "Get a unique Request ID (e.g., #REQ-2026-8941).",
                 "Watch live 4-step progress bar on your phone screen."
             ],
             badge="LIVE TRACKING")

    add_card(s7, Inches(9.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Step 4: Resolved!",
             [
                 "Officer inspects and completes the work.",
                 "Status updates to 'Resolved'.",
                 "Download official completion receipt."
             ],
             badge="COMPLETION")

    # ==========================================
    # SLIDE 8: OFFICER INSPECTION & RESOLUTION
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, "Officer Review & Inspection Workflow", "OFFICER EXPERIENCE")

    add_card(s8, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "How Officers Manage Complaints",
             [
                 "Ward Queue: Officers see all incoming citizen requests organized neatly.",
                 "Photo & Map Inspection: Check uploaded photos and GPS location before visiting.",
                 "Status Upgrades: Move status: 'Submitted' ➔ 'Under Review' ➔ 'Action Taken'.",
                 "Resolution Notes: Add comments explaining the repair work done.",
                 "Mark as Complete: One-click closure when work is verified."
             ],
             badge="INSPECTION PIPELINE")

    add_card(s8, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "Why Officers Love It",
             [
                 "No Paper Files: Zero paperwork clutter on office desks.",
                 "Priority Sorting: High-priority emergency issues appear on top.",
                 "Clear Proof: Citizen photos prevent confusion about the exact defect.",
                 "Audit Records: Clear proof of work done by the officer and team."
             ],
             badge="OFFICER BENEFITS")

    # ==========================================
    # SLIDE 9: SUPER ADMIN MASTER CMS
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, "Super Administrator Master CMS Portal", "ADMINISTRATION")

    add_card(s9, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "Full Control & Management",
             [
                 "User Accounts: Add, edit, or deactivate officer and staff accounts.",
                 "City Corporations CMS: Manage 5 City Corporations and 198 BBMP Wards.",
                 "Department Directory: Update department contacts, head officers, and emails.",
                 "Public Notices & Alerts: Publish city news, weather warnings, and gazettes.",
                 "Events Management: Add town hall meetings and public events to calendar."
             ],
             badge="MASTER CMS")

    add_card(s9, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "Real-Time City Analytics",
             [
                 "Track total complaints submitted vs. resolved across all wards.",
                 "Identify wards needing more infrastructure attention.",
                 "Monitor average resolution time by department.",
                 "Export performance data for council meetings."
             ],
             badge="DATA & INSIGHTS")

    # ==========================================
    # SLIDE 10: INTERACTIVE MAP & PUBLIC FEATURES
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, "Interactive City Map & Public Features", "CITIZEN TOOLS")

    add_card(s10, Inches(0.8), Inches(1.75), Inches(3.6), Inches(5.1),
             "Interactive GIS Landmark Map",
             [
                 "Visual city map pinpointing key municipal landmarks.",
                 "Pin overlays for BBMP Head Office, Town Hall, Clinics & Parks.",
                 "RHS 'Get Directions' button opening navigation route.",
                 "Touch-friendly for mobile maps."
             ],
             badge="MAPS & ROUTING")

    add_card(s10, Inches(4.85), Inches(1.75), Inches(3.6), Inches(5.1),
             "City Directory & Events",
             [
                 "Searchable directory of all ward officers and departments.",
                 "Calendar of public council hearings and festival dates.",
                 "Official municipal bylaws and application forms library.",
                 "Photo gallery of city development milestones."
             ],
             badge="PUBLIC INFORMATION")

    add_card(s10, Inches(8.9), Inches(1.75), Inches(3.6), Inches(5.1),
             "24/7 Helpline & Emergency",
             [
                 "BBMP Sahaya 24/7 Control Room (1533).",
                 "Quick 1-tap call button in mobile menu drawer.",
                 "Live alert marquee showing urgent city advisories.",
                 "Grievance escalation cell for pending issues."
             ],
             badge="EMERGENCY SUPPORT")

    # ==========================================
    # SLIDE 11: TECHNOLOGY EXPLAINED SIMPLY
    # ==========================================
    s11 = prs.slides.add_slide(blank_layout)
    add_header(s11, "The Technology Stack (Explained Simply)", "SYSTEM ARCHITECTURE")

    add_card(s11, Inches(0.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Frontend",
             [
                 "React 18 & Vite.",
                 "Builds fast, interactive web screens.",
                 "No page reloads needed.",
                 "Modern icons & clean visuals."
             ],
             badge="USER INTERFACE")

    add_card(s11, Inches(3.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Backend",
             [
                 "Python & Django REST.",
                 "Processes citizen applications.",
                 "Enforces security & user roles.",
                 "Fast serverless API."
             ],
             badge="LOGIC & API")

    add_card(s11, Inches(6.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Cloud Database",
             [
                 "Supabase PostgreSQL.",
                 "Stores all users, complaints, and records safely in cloud.",
                 "Zero data loss risk."
             ],
             badge="DATABASE")

    add_card(s11, Inches(9.8), Inches(1.75), Inches(2.7), Inches(5.1),
             "Media CDN",
             [
                 "Cloudinary CDN.",
                 "Stores all photo evidence & documents.",
                 "Loads images instantly on mobile and desktop."
             ],
             badge="MEDIA STORAGE")

    # ==========================================
    # SLIDE 12: MOBILE, TABLET & LAPTOP DESIGN
    # ==========================================
    s12 = prs.slides.add_slide(blank_layout)
    add_header(s12, "Designed for Mobile, Tablet & Laptop Screens", "RESPONSIVE UI/UX")

    add_card(s12, Inches(0.8), Inches(1.75), Inches(3.6), Inches(5.1),
             "📱 Mobile Phone View (< 768px)",
             [
                 "Neat single-column vertical cards.",
                 "Big, readable 17.5px text for easy reading.",
                 "Slide-out mobile navigation drawer with active user status.",
                 "Large 48px touch buttons for easy tapping.",
                 "Smooth horizontal scroll for tables."
             ],
             badge="MOBILE PHONES")

    add_card(s12, Inches(4.85), Inches(1.75), Inches(3.6), Inches(5.1),
             "📱 Tablet View (768px – 1024px)",
             [
                 "Clean 2-column card layouts.",
                 "Hamburger drawer toggle.",
                 "Scrollable category filter chips.",
                 "Touch-optimized forms and modals."
             ],
             badge="TABLETS / IPADS")

    add_card(s12, Inches(8.9), Inches(1.75), Inches(3.6), Inches(5.1),
             "💻 Laptop / Desktop View (> 1024px)",
             [
                 "Full widescreen 4-column card grid.",
                 "Fixed sticky citizen dashboard navigation bar.",
                 "Expanded top menu bar with sub-dropdowns.",
                 "Comprehensive administrative data tables."
             ],
             badge="LAPTOPS & COMPUTERS")

    # ==========================================
    # SLIDE 13: SECURITY & DATA PRIVACY
    # ==========================================
    s13 = prs.slides.add_slide(blank_layout)
    add_header(s13, "Security & Data Privacy (Keeping Records Safe)", "SECURITY ARCHITECTURE")

    add_card(s13, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.1),
             "Data Protection & Privacy",
             [
                 "Password Hashing: Citizen and officer passwords are securely encrypted.",
                 "Role Protection: Only authorized officers can approve requests in their ward.",
                 "Zero Secret Leakage: All database passwords & keys are stored in secure .env files.",
                 "Never stored in public GitHub files."
             ],
             badge="PROTECTION")

    add_card(s13, Inches(6.9), Inches(1.75), Inches(5.6), Inches(5.1),
             "Cloud Reliability & Backups",
             [
                 "Automated Cloud Backups: Database backed up continuously on Supabase.",
                 "Secure HTTPS: All traffic encrypted with modern SSL certificates.",
                 "Audit Logs: Records who approved which certificate and when.",
                 "High Availability: 99.9% uptime powered by Vercel and Supabase cloud."
             ],
             badge="RELIABILITY")

    # ==========================================
    # SLIDE 14: FUTURE ROADMAP & INNOVATIONS
    # ==========================================
    s14 = prs.slides.add_slide(blank_layout)
    add_header(s14, "Future Roadmap & Exciting Possibilities", "FUTURE INNOVATION")

    add_card(s14, Inches(0.8), Inches(1.75), Inches(3.6), Inches(5.1),
             "🤖 AI Pothole Detection",
             [
                 "AI analyzes citizen photos automatically.",
                 "Measures pothole depth and hazard level.",
                 "Auto-prioritizes dangerous road hazards for immediate repair."
             ],
             badge="AI INTEGRATION")

    add_card(s14, Inches(4.85), Inches(1.75), Inches(3.6), Inches(5.1),
             "💬 WhatsApp Civic Bot",
             [
                 "File complaints directly via WhatsApp message.",
                 "Send photo and location in WhatsApp chat.",
                 "Instant status check by typing Request ID."
             ],
             badge="WHATSAPP BOT")

    add_card(s14, Inches(8.9), Inches(1.75), Inches(3.6), Inches(5.1),
             "🔔 SMS & Voice Broadcasts",
             [
                 "Instant SMS alert when complaint is resolved.",
                 "Voice call alerts for emergency flood or weather warnings.",
                 "Local language translation (Kannada / Hindi / English)."
             ],
             badge="COMMUNICATION")

    # ==========================================
    # SLIDE 15: CONCLUSION & LIVE LINKS
    # ==========================================
    s15 = prs.slides.add_slide(blank_layout)
    bg15 = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg15.fill.solid()
    bg15.fill.fore_color.rgb = COLOR_NAVY
    bg15.line.fill.background()

    strip15 = s15.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.35), Inches(7.5))
    strip15.fill.solid()
    strip15.fill.fore_color.rgb = COLOR_CYAN
    strip15.line.fill.background()

    tb15 = s15.shapes.add_textbox(Inches(1.2), Inches(1.2), Inches(11.0), Inches(5.0))
    tf15 = tb15.text_frame
    tf15.word_wrap = True

    p15_1 = tf15.paragraphs[0]
    p15_1.text = "THANK YOU!"
    p15_1.font.size = Pt(38)
    p15_1.font.bold = True
    p15_1.font.color.rgb = COLOR_WHITE
    p15_1.space_after = Pt(12)

    p15_2 = tf15.add_paragraph()
    p15_2.text = "The Municipal Automation Portal delivers a smart, transparent, and hassle-free civic governance experience for everyone."
    p15_2.font.size = Pt(16)
    p15_2.font.color.rgb = COLOR_LIGHT_CYAN
    p15_2.space_after = Pt(24)

    p15_3 = tf15.add_paragraph()
    p15_3.text = "🌐 Live Application Links:\n•  Frontend Web App:  https://municipality-automation-frontend.vercel.app/\n•  Backend REST API:  https://municipality-automation-backend.vercel.app/\n•  GitHub Repository:  https://github.com/pratiksha-harode005/municipality_automation.git"
    p15_3.font.size = Pt(13)
    p15_3.font.color.rgb = RGBColor(226, 232, 240)
    p15_3.space_after = Pt(24)

    p15_4 = tf15.add_paragraph()
    p15_4.text = "Questions, Feedback & Discussion Welcome!"
    p15_4.font.size = Pt(16)
    p15_4.font.bold = True
    p15_4.font.color.rgb = COLOR_GOLD

    output_path = os.path.abspath("Municipality_Automation_Project_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation with 15 slides saved successfully at: {output_path}")

if __name__ == "__main__":
    create_deck()
